from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, List, Optional

from pptx import Presentation
from pptx.text.text import _Paragraph

from .utils import (
    PARAGRAPH_DELIMITER,
    StylePalette,
    chunked,
    escape_text,
    iter_tagged_segments,
    unescape_text,
)

PPTX_SYSTEM_PROMPT = (
    "You are a meticulous translator. Translate the text into English while preserving the <sX> tags exactly. "
    f"Paragraphs are separated with {PARAGRAPH_DELIMITER}."
)


@dataclass(slots=True)
class PptxParagraphBundle:
    paragraph: _Paragraph
    palette: StylePalette
    prompt: str


def iter_pptx_paragraphs(presentation: Presentation) -> Iterable[_Paragraph]:
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                yield paragraph


def build_paragraph_bundle(paragraph: _Paragraph) -> Optional[PptxParagraphBundle]:
    palette = StylePalette()
    segments: List[str] = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        style_key = palette.register_pptx_run(run)
        segments.append(f"<{style_key}>{escape_text(text)}</{style_key}>")
    prompt = "".join(segments)
    if not prompt:
        return None
    return PptxParagraphBundle(paragraph=paragraph, palette=palette, prompt=prompt)


def clear_paragraph(paragraph: _Paragraph) -> None:
    for run in list(paragraph.runs):
        paragraph._p.remove(run._r)


def apply_translation(bundle: PptxParagraphBundle, translated_payload: str) -> None:
    segments = iter_tagged_segments(translated_payload)
    if not segments:
        return
    clear_paragraph(bundle.paragraph)
    for key, text in segments:
        run = bundle.paragraph.add_run()
        run.text = unescape_text(text)
        if bundle.palette.has(key):
            bundle.palette.apply_pptx(key, run)


async def translate_pptx_file(
    input_path: Path,
    output_path: Path,
    llm_client,
    progress_callback: Callable[[float, str], None] | None = None,
    batch_size: int = 8,
) -> None:
    presentation = Presentation(input_path)
    bundles = [bundle for paragraph in iter_pptx_paragraphs(presentation) if (bundle := build_paragraph_bundle(paragraph))]

    if not bundles:
        presentation.save(output_path)
        return

    total = len(bundles)
    processed = 0

    for chunk in chunked(bundles, batch_size):
        prompts = [item.prompt for item in chunk]
        user_prompt = PARAGRAPH_DELIMITER.join(prompts)
        translated = await llm_client.translate(system_prompt=PPTX_SYSTEM_PROMPT, user_prompt=user_prompt)
        responses = translated.split(PARAGRAPH_DELIMITER)
        if len(responses) != len(prompts):
            responses = (responses + [responses[-1]] * len(prompts))[: len(prompts)] if responses else [translated] * len(prompts)
        for bundle, payload in zip(chunk, responses):
            apply_translation(bundle, payload)
            processed += 1
            if progress_callback:
                progress_callback(processed / total, f"Translated slide text {processed}/{total}")

    presentation.save(output_path)
